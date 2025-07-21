# document_processor.py
import re
import pickle
import os
import base64
from io import BytesIO
import numpy as np
import sqlite3
import nltk
from nltk.tokenize import sent_tokenize
from sentence_transformers import SentenceTransformer
import fitz
import streamlit as st
from typing import Dict, List, Optional
from docx import Document
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import json

# Download required NLTK data
try:
    nltk.data.find('tokenizers/punkt')
except LookupError:
    nltk.download('punkt')
try:
    nltk.data.find('tokenizers/punkt_tab')
except LookupError:
    try:
        nltk.download('punkt_tab')
    except:
        pass

class DocumentProcessor:
    """Handles document ingestion and processing from uploads and project folders"""
    
    def __init__(self, db_manager, documents_folder: str = "documents"):
        self.db_manager = db_manager
        self.model = SentenceTransformer('all-MiniLM-L6-v2')
        self.documents_folder = documents_folder
        
        # Create documents folder if it doesn't exist
        if not os.path.exists(self.documents_folder):
            os.makedirs(self.documents_folder)
    
    def load_documents_from_folder(self) -> bool:
        """Load all documents from the project documents folder"""
        if not os.path.exists(self.documents_folder):
            st.warning(f"Documents folder '{self.documents_folder}' not found")
            return False
        
        processed_count = 0
        error_count = 0
        
        # Get list of already processed documents
        processed_docs = self._get_processed_documents()
        
        for filename in os.listdir(self.documents_folder):
            file_path = os.path.join(self.documents_folder, filename)
            
            # Skip if already processed
            if filename in processed_docs:
                continue
                
            if os.path.isfile(file_path):
                file_extension = filename.lower().split('.')[-1]
                
                if file_extension in ['pdf', 'docx', 'pptx', 'ppt']:
                    try:
                        with open(file_path, 'rb') as file:
                            file_content = file.read()
                        
                        if self.ingest_document(file_content, filename, file_extension):
                            processed_count += 1
                            st.success(f"Processed: {filename}")
                        else:
                            error_count += 1
                            st.error(f"Failed to process: {filename}")
                            
                    except Exception as e:
                        error_count += 1
                        st.error(f"Error reading {filename}: {str(e)}")
        
        if processed_count > 0:
            st.success(f"Successfully processed {processed_count} documents from folder")
        if error_count > 0:
            st.warning(f"Failed to process {error_count} documents")
            
        return processed_count > 0
    
    def _get_processed_documents(self) -> set:
        """Get list of already processed document filenames"""
        try:
            conn = sqlite3.connect(self.db_manager.db_path)
            cursor = conn.cursor()
            cursor.execute("SELECT filename FROM documents")
            processed = {row[0] for row in cursor.fetchall()}
            conn.close()
            return processed
        except Exception:
            return set()
    
    def process_pdf(self, file_content: bytes, filename: str) -> Dict:
        """Extract text, tables, images, and structured content from PDF with improved parsing"""
        text_content = ""
        images = []
        tables = []
        structured_content = []
        
        try:
            pdf_doc = fitz.open(stream=file_content, filetype="pdf")
            
            for page_num in range(len(pdf_doc)):
                page = pdf_doc.load_page(page_num)
                
                # Extract raw text
                page_text = page.get_text()
                
                # Extract structured text with formatting
                text_dict = page.get_text("dict")
                structured_page_content = self._extract_structured_content_from_page(text_dict, page_num + 1)
                structured_content.extend(structured_page_content)
                
                # Extract tables using text blocks analysis
                page_tables = self._extract_tables_from_page(page, page_num + 1)
                tables.extend(page_tables)
                
                # Build comprehensive text content
                page_content = f"\n=== PAGE {page_num + 1} ===\n"
                page_content += page_text
                
                # Add structured table content
                if page_tables:
                    page_content += "\n--- TABLES ---\n"
                    for table in page_tables:
                        page_content += table['formatted_text'] + "\n"
                
                text_content += page_content + "\n"
                
                # Extract images with enhanced context association
                image_list = page.get_images()
                for img_index, img in enumerate(image_list):
                    try:
                        xref = img[0]
                        pix = fitz.Pixmap(pdf_doc, xref)
                        if pix.n < 5:  # GRAY or RGB
                            img_data = pix.tobytes("png")
                            img_b64 = base64.b64encode(img_data).decode('utf-8')
                            
                            # Enhanced context analysis for error codes and table associations
                            context_text = self._get_enhanced_image_context(pdf_doc, page_num, page_text, page_tables)
                            associated_codes = self.find_error_codes_on_page(context_text)
                            
                            # Try to associate with nearby tables
                            associated_tables = [t['id'] for t in page_tables]
                            
                            images.append({
                                "data": img_b64,
                                "page": page_num + 1,
                                "index": img_index,
                                "associated_codes": associated_codes,
                                "associated_tables": associated_tables,
                                "page_text": page_text[:500],
                                "context": context_text[:300]
                            })
                        pix = None
                    except Exception as e:
                        print(f"Error processing image {img_index} on page {page_num}: {e}")
                        continue
            
            pdf_doc.close()
        except Exception as e:
            st.error(f"Error processing PDF {filename}: {str(e)}")
            return {"text": "", "images": [], "tables": [], "structured_content": [], "error": str(e)}
        
        return {
            "text": text_content, 
            "images": images, 
            "tables": tables,
            "structured_content": structured_content
        }

    def _extract_structured_content_from_page(self, text_dict: dict, page_num: int) -> List[Dict]:
        """Extract structured content like headings, lists, and formatted sections"""
        structured_items = []
        
        try:
            for block in text_dict.get("blocks", []):
                if "lines" in block:
                    # Analyze text blocks for structure
                    block_text = ""
                    font_sizes = []
                    is_bold = False
                    
                    for line in block["lines"]:
                        for span in line.get("spans", []):
                            span_text = span.get("text", "").strip()
                            if span_text:
                                block_text += span_text + " "
                                font_sizes.append(span.get("size", 12))
                                if span.get("flags", 0) & 2**4:  # Bold flag
                                    is_bold = True
                    
                    block_text = block_text.strip()
                    if not block_text:
                        continue
                    
                    # Determine content type
                    avg_font_size = sum(font_sizes) / len(font_sizes) if font_sizes else 12
                    content_type = self._classify_content_type(block_text, avg_font_size, is_bold)
                    
                    if content_type != "regular_text":
                        structured_items.append({
                            "type": content_type,
                            "text": block_text,
                            "page": page_num,
                            "font_size": avg_font_size,
                            "is_bold": is_bold
                        })
        
        except Exception as e:
            print(f"Error extracting structured content from page {page_num}: {e}")
        
        return structured_items

    def _classify_content_type(self, text: str, font_size: float, is_bold: bool) -> str:
        """Classify text content type based on formatting and patterns"""
        text_lower = text.lower().strip()
        
        # Headers/Titles (larger font, bold, or specific patterns)
        if (font_size > 14 or is_bold) and len(text) < 100:
            if any(keyword in text_lower for keyword in ['code', 'error', 'fault', 'diagnostic']):
                return "section_header"
            if text.isupper() or (is_bold and len(text.split()) <= 8):
                return "title"
        
        # Lists (numbered or bulleted)
        if re.match(r'^\d+\.?\s', text) or re.match(r'^[•\-\*]\s', text):
            return "list_item"
        
        # Table headers or structured data
        if '|' in text or '\t' in text or (len(text.split()) <= 10 and any(c.isdigit() for c in text)):
            return "tabular_data"
        
        # Error codes or technical references
        if re.search(r'\b[A-F0-9]{4}\b', text) or any(keyword in text_lower for keyword in ['msk', 'blt', 'code']):
            return "error_reference"
        
        # Steps or procedures
        if re.match(r'^\d+\.\s', text) and len(text) > 20:
            return "procedure_step"
        
        return "regular_text"

    def _extract_tables_from_page(self, page, page_num: int) -> List[Dict]:
        """Extract and format tables from PDF page using text block analysis"""
        tables = []
        
        try:
            # Get text as dictionary with positioning info
            text_dict = page.get_text("dict")
            
            # Try to find table-like structures
            table_blocks = []
            for block in text_dict.get("blocks", []):
                if "lines" in block:
                    block_text = ""
                    for line in block["lines"]:
                        line_text = ""
                        for span in line.get("spans", []):
                            line_text += span.get("text", "")
                        block_text += line_text + "\n"
                    
                    # Check if block looks like a table
                    if self._is_table_block(block_text):
                        table_blocks.append({
                            "text": block_text,
                            "bbox": block.get("bbox", [0, 0, 0, 0])
                        })
            
            # Process identified table blocks
            for i, table_block in enumerate(table_blocks):
                formatted_table = self._format_table_text(table_block["text"])
                if formatted_table:
                    tables.append({
                        "id": f"page_{page_num}_table_{i}",
                        "page": page_num,
                        "raw_text": table_block["text"],
                        "formatted_text": formatted_table,
                        "bbox": table_block["bbox"]
                    })
        
        except Exception as e:
            print(f"Error extracting tables from page {page_num}: {e}")
        
        return tables

    def _is_table_block(self, text: str) -> bool:
        """Determine if a text block represents a table"""
        lines = text.strip().split('\n')
        if len(lines) < 2:
            return False
        
        # Check for table indicators
        table_indicators = 0
        
        # Look for consistent column separators
        if any('|' in line for line in lines):
            table_indicators += 2
        
        # Look for consistent spacing patterns
        if len(set(len(line) for line in lines if line.strip())) <= 3:
            table_indicators += 1
        
        # Look for headers with data rows
        first_line = lines[0].strip()
        if any(keyword in first_line.lower() for keyword in ['position', 'code', 'description', 'dérangements']):
            table_indicators += 2
        
        # Look for structured data patterns (numbers, codes)
        structured_lines = sum(1 for line in lines if re.search(r'\d+', line) or re.search(r'[XxOo-]{2,}', line))
        if structured_lines > len(lines) * 0.5:
            table_indicators += 1
        
        return table_indicators >= 2

    def _format_table_text(self, raw_table_text: str) -> str:
        """Format raw table text into a more structured format"""
        lines = [line.strip() for line in raw_table_text.strip().split('\n') if line.strip()]
        if not lines:
            return ""
        
        formatted_lines = []
        
        # Try to identify header and data rows
        for i, line in enumerate(lines):
            # Clean up the line
            clean_line = re.sub(r'\s+', ' ', line)
            
            # Add visual separation for what looks like headers
            if i == 0 or any(keyword in line.lower() for keyword in ['position', 'code', 'description']):
                formatted_lines.append(f"| {clean_line} |")
                formatted_lines.append("|" + "-" * (len(clean_line) + 2) + "|")
            else:
                formatted_lines.append(f"| {clean_line} |")
        
        return '\n'.join(formatted_lines)
    
    def process_pptx(self, file_content: bytes, filename: str) -> Dict:
        """Extract text and images from PowerPoint presentation"""
        text_content = ""
        images = []
        
        try:
            prs = Presentation(BytesIO(file_content))
            
            for slide_num, slide in enumerate(prs.slides):
                slide_text = ""
                
                # Extract text from all shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text += shape.text + "\n"
                    
                    # Extract images from shapes
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        try:
                            image = shape.image
                            img_data = image.blob
                            img_b64 = base64.b64encode(img_data).decode('utf-8')
                            
                            # Try to associate with error codes found on this slide
                            associated_codes = self.find_error_codes_on_page(slide_text)
                            
                            images.append({
                                "data": img_b64,
                                "page": slide_num + 1,  # Use slide number as page
                                "index": len(images),
                                "associated_codes": associated_codes,
                                "page_text": slide_text[:500]
                            })
                        except Exception as e:
                            print(f"Error processing image on slide {slide_num}: {e}")
                            continue
                
                # Add slide header
                text_content += f"=== Slide {slide_num + 1} ===\n"
                text_content += slide_text + "\n\n"
                
                # Extract text from notes
                if slide.notes_slide and slide.notes_slide.notes_text_frame:
                    notes_text = slide.notes_slide.notes_text_frame.text
                    if notes_text.strip():
                        text_content += f"Notes: {notes_text}\n\n"
        
        except Exception as e:
            st.error(f"Error processing PowerPoint {filename}: {str(e)}")
            return {"text": "", "images": [], "error": str(e)}
        
        return {"text": text_content, "images": images}
    
    def find_error_codes_on_page(self, text: str) -> List[str]:
        """Find error codes mentioned on a specific page with improved detection"""
        # Multiple patterns for better error code detection
        patterns = [
            r'\b([A-Fa-f0-9]{4})\b',  # Basic hex pattern
            r'Code\s+([A-Fa-f0-9]{4})',  # Code XXXX
            r'Error\s+([A-Fa-f0-9]{4})',  # Error XXXX
            r'([A-Fa-f0-9]{4})\s*[-:]\s*',  # XXXX: or XXXX-
        ]
        
        codes = set()
        for pattern in patterns:
            matches = re.findall(pattern, text.upper())
            codes.update(matches)
        
        return list(codes)
    
    def process_docx(self, file_content: bytes, filename: str) -> Dict:
        """Extract text from Word document"""
        try:
            doc = Document(BytesIO(file_content))
            text_content = ""
            
            for paragraph in doc.paragraphs:
                text_content += paragraph.text + "\n"
            
            # Extract tables
            tables_text = ""
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        row_text.append(cell.text.strip())
                    tables_text += " | ".join(row_text) + "\n"
            
            return {"text": text_content + "\n" + tables_text, "images": []}
        except Exception as e:
            st.error(f"Error processing DOCX {filename}: {str(e)}")
            return {"text": "", "images": [], "error": str(e)}
    
    def extract_error_codes_with_context(self, text: str) -> List[Dict]:
        """Enhanced error code extraction with better pattern recognition for technical documents"""
        error_codes = []
        
        # Enhanced patterns for different error code formats
        patterns = [
            # Standard hex codes with descriptions
            r'([A-Fa-f0-9]{4})\s*[-:=]\s*(.{10,300})',
            # Codes in tables or structured format
            r'Code\s+([A-Fa-f0-9]{4})\s*[-:]?\s*(.{10,300})',
            r'Error\s+([A-Fa-f0-9]{4})\s*[-:]?\s*(.{10,300})',
            # Technical codes (MSK, BLT, etc.) followed by descriptions
            r'(MSK[0-9]|BLT|UVM)\s*[-:]?\s*(.{20,300})',
            # Position-based codes (like in your table)
            r'([0-9]+\.?)\s+([XxOo\s-]{8,})\s+(.{20,300})',
            # French technical terms with codes
            r'([A-Fa-f0-9]{4})\s+(.{20,300}?)(?:température|moteur|circuit|défaut)',
        ]
        
        lines = text.split('\n')
        for i, line in enumerate(lines):
            line = line.strip()
            if not line or len(line) < 10:
                continue
            
            for pattern in patterns:
                matches = re.findall(pattern, line, re.IGNORECASE | re.MULTILINE)
                for match in matches:
                    if len(match) >= 2:
                        code = match[0].upper().strip()
                        description = match[1].strip()
                        
                        # Clean up description
                        description = re.sub(r'\s+', ' ', description)
                        description = description[:300]  # Limit length
                        
                        # Skip if description is too short or just symbols
                        if len(description) < 15 or re.match(r'^[XxOo\s-]+$', description):
                            continue
                        
                        # Get extended context (surrounding lines)
                        context_lines = []
                        for j in range(max(0, i-3), min(len(lines), i+4)):
                            if lines[j].strip() and j != i:
                                context_lines.append(lines[j].strip())
                        
                        # Check for duplicate codes and merge if similar
                        existing_code = next((ec for ec in error_codes if ec['code'] == code), None)
                        if existing_code:
                            # Merge descriptions if different
                            if description not in existing_code['description']:
                                existing_code['description'] += f" | {description}"
                                existing_code['context'] += f" | {' '.join(context_lines)}"
                        else:
                            error_codes.append({
                                "code": code,
                                "description": description,
                                "context": ' '.join(context_lines),
                                "line_number": i + 1,
                                "source_type": "technical_document"
                            })
        
        return error_codes

    def _get_enhanced_image_context(self, pdf_doc, page_num: int, page_text: str, page_tables: List) -> str:
        """Get enhanced context for image including surrounding pages and table content"""
        context_parts = []
        
        # Current page text
        context_parts.append(page_text)
        
        # Previous page context
        if page_num > 0:
            try:
                prev_page = pdf_doc.load_page(page_num - 1)
                prev_text = prev_page.get_text()
                context_parts.append(prev_text[-300:])  # Last 300 chars
            except:
                pass
        
        # Next page context
        if page_num < len(pdf_doc) - 1:
            try:
                next_page = pdf_doc.load_page(page_num + 1)
                next_text = next_page.get_text()
                context_parts.append(next_text[:300])  # First 300 chars
            except:
                pass
        
        # Table context from current page
        table_context = []
        for table in page_tables:
            table_context.append(table['formatted_text'])
        context_parts.extend(table_context)
        
        return "\n\n".join(context_parts)
    
    def create_embeddings(self, texts: List[str]) -> np.ndarray:
        """Create vector embeddings for text chunks"""
        return self.model.encode(texts)
    
    def ingest_document(self, file_content: bytes, filename: str, file_type: str) -> bool:
        """Enhanced document ingestion with support for tables and structured content"""
        conn = None
        try:
            # Process document based on type
            if file_type == "pdf":
                result = self.process_pdf(file_content, filename)
            elif file_type == "docx":
                result = self.process_docx(file_content, filename)
            elif file_type == "pptx":
                result = self.process_pptx(file_content, filename)
            else:
                st.error(f"Unsupported file type: {file_type}")
                return False
            
            if "error" in result:
                return False
            
            text_content = result["text"]
            images = result.get("images", [])
            tables = result.get("tables", [])
            structured_content = result.get("structured_content", [])
            
            # Store document
            conn = sqlite3.connect(self.db_manager.db_path, timeout=30.0)
            cursor = conn.cursor()
            
            conn.execute("BEGIN IMMEDIATE")
            
            # Enhanced metadata
            metadata = {
                "image_count": len(images),
                "table_count": len(tables),
                "structured_sections": len(structured_content)
            }
            if file_type == "pptx":
                slide_count = text_content.count("=== Slide")
                metadata["slide_count"] = slide_count
            
            cursor.execute("""
                INSERT INTO documents (filename, content, doc_type, metadata) 
                VALUES (?, ?, ?, ?)
            """, (filename, text_content, file_type, json.dumps(metadata)))
            
            doc_id = cursor.lastrowid
            
            # Store tables separately for better search
            for table in tables:
                cursor.execute("""
                    INSERT INTO tables (doc_id, table_id, page_number, raw_content, formatted_content, bbox)
                    VALUES (?, ?, ?, ?, ?, ?)
                """, (
                    doc_id,
                    table["id"],
                    table["page"],
                    table["raw_text"],
                    table["formatted_text"],
                    json.dumps(table["bbox"])
                ))
            
            # Enhanced error code extraction
            error_codes = self.extract_error_codes_with_context(text_content)
            
            # Also extract from table content
            for table in tables:
                table_error_codes = self.extract_error_codes_with_context(table["raw_text"])
                for code in table_error_codes:
                    code["source_type"] = "table"
                    code["table_id"] = table["id"]
                error_codes.extend(table_error_codes)
            
            # Store error codes
            # Store error codes (FIXED COLUMN LIST)
            stored_error_codes = {}
            for error_code in error_codes:
                cursor.execute("""
                    INSERT OR REPLACE INTO error_codes 
                    (code, description, source_doc, procedure_steps, source_type, metadata) 
                    VALUES (?, ?, ?, ?, ?, ?)
                """, (
                    error_code["code"], 
                    error_code["description"], 
                    filename,
                    error_code["context"],
                    error_code.get("source_type", "document"),
                    json.dumps({
                        "line_number": error_code.get("line_number"),
                        "table_id": error_code.get("table_id")
                    })
                ))
                stored_error_codes[error_code["code"]] = error_code
            
            # Store images with enhanced associations
            for i, img in enumerate(images):
                if file_type == "pptx":
                    img_filename = f"{filename}_slide_{img['page']}_img_{i}"
                else:
                    img_filename = f"{filename}_page_{img['page']}_img_{i}"
                
                # Enhanced error code association
                best_associated_code = None
                if img.get("associated_codes"):
                    for code in img["associated_codes"]:
                        if code in stored_error_codes:
                            best_associated_code = code
                            break
                    if not best_associated_code and img["associated_codes"]:
                        best_associated_code = img["associated_codes"][0]
                
                # Enhanced description with table and context information
                enhanced_description = f"Image from {filename} {'slide' if file_type == 'pptx' else 'page'} {img['page']}"
                
                # Add context from nearby tables
                if img.get('associated_tables'):
                    table_contexts = [t['formatted_text'][:100] for t in tables 
                                    if t['id'] in img['associated_tables']]
                    if table_contexts:
                        enhanced_description += f" - Table context: {' | '.join(table_contexts)}"
                
                # Add general context
                if img.get('context'):
                    context_clean = re.sub(r'\s+', ' ', img['context'][:200])
                    enhanced_description += f" - Context: {context_clean}"
                
                # Store image
                cursor.execute("""
                    INSERT INTO images 
                    (filename, image_data, description, step_number, associated_code, metadata) 
                    VALUES (?, ?, ?, ?, ?, ?)
                """, (
                    img_filename, 
                    img["data"],
                    enhanced_description,
                    img['page'],
                    best_associated_code,
                    json.dumps({
                        "associated_tables": img.get('associated_tables', []),
                        "context_length": len(img.get('context', ''))
                    })
                ))
            
            # Create embeddings for meaningful content
            all_content_for_embeddings = []
            
            # Regular sentences
            try:
                sentences = sent_tokenize(text_content)
            except LookupError:
                sentences = text_content.split('. ')
                sentences = [s.strip() + '.' for s in sentences if s.strip()]
            
            meaningful_sentences = [
                s for s in sentences 
                if len(s.strip()) > 30 and  
                not s.strip().startswith("=== ") and
                not s.strip().startswith("Page") and
                len(s.split()) > 5
            ]
            all_content_for_embeddings.extend(meaningful_sentences)
            
            # Table content for embeddings
            for table in tables:
                table_lines = [line.strip() for line in table["formatted_text"].split('\n') 
                            if line.strip() and not line.startswith('|---')]
                all_content_for_embeddings.extend(table_lines)
            
            # Structured content for embeddings
            for struct_item in structured_content:
                if len(struct_item["text"]) > 20:
                    all_content_for_embeddings.append(struct_item["text"])
            
            # Create embeddings in batches
            if all_content_for_embeddings:
                batch_size = 50
                for i in range(0, len(all_content_for_embeddings), batch_size):
                    batch = all_content_for_embeddings[i:i+batch_size]
                    embeddings = self.create_embeddings(batch)
                    
                    for content, embedding in zip(batch, embeddings):
                        cursor.execute("""
                            INSERT INTO embeddings (content_id, content_type, embedding, text_content) 
                            VALUES (?, ?, ?, ?)
                        """, (str(doc_id), "document", pickle.dumps(embedding), content))
            
            conn.commit()
            
            # Print comprehensive summary
            summary_parts = [f"Document processed: {filename}"]
            if images:
                summary_parts.append(f"{len(images)} images")
            if tables:
                summary_parts.append(f"{len(tables)} tables")
            if error_codes:
                summary_parts.append(f"{len(error_codes)} error codes")
            
            st.success(" | ".join(summary_parts))
            
            return True
            
        except Exception as e:
            if conn:
                conn.rollback()
            st.error(f"Error ingesting document {filename}: {str(e)}")
            print(f"Detailed error: {e}")
            return False
        finally:
            if conn:
                conn.close()
    
    def get_document_stats(self) -> Dict:
        """Get statistics about processed documents"""
        try:
            conn = sqlite3.connect(self.db_manager.db_path)
            cursor = conn.cursor()
            
            cursor.execute("SELECT COUNT(*) FROM documents")
            doc_count = cursor.fetchone()[0]
            
            cursor.execute("SELECT COUNT(*) FROM error_codes")
            error_code_count = cursor.fetchone()[0]
            
            cursor.execute("SELECT COUNT(*) FROM images")
            image_count = cursor.fetchone()[0]
            
            cursor.execute("SELECT COUNT(*) FROM embeddings")
            embedding_count = cursor.fetchone()[0]
            
            # Get document type breakdown
            cursor.execute("SELECT doc_type, COUNT(*) FROM documents GROUP BY doc_type")
            doc_types = dict(cursor.fetchall())
            
            conn.close()
            
            return {
                "documents": doc_count,
                "error_codes": error_code_count,
                "images": image_count,
                "embeddings": embedding_count,
                "document_types": doc_types
            }
        except Exception:
            return {
                "documents": 0, 
                "error_codes": 0, 
                "images": 0, 
                "embeddings": 0,
                "document_types": {}
            }
